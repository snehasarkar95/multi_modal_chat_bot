import PyPDF2
import docx
import pandas as pd
from PIL import Image
import io
import json
import xml.etree.ElementTree as ET
import yaml
import markdown
import logging
from pptx import Presentation
import regex as re
from pathlib import Path
import mimetypes
import magic

def process_uploaded_file(file_path, filename, file_ext):
    """Process uploaded files based on MIME type or extension"""
    result = {"type": "unknown", "metadata": {}}
    sheets_data = []
    preview = None
    file_ext = file_ext.lower() if file_ext else ""
    mime_type = mime_type.lower() if mime_type else ""

    try:
        # Text Documents
        if (file_ext in ["txt", "log"] or 
            mime_type.startswith("text/") or 
            mime_type in ["application/x-log", "application/log"]):
            
            # Try different encodings to handle various text files
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = ""
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            result.update({
                "type": "text",
                "content": content,
                "metadata": {
                    "length": len(content), 
                    "lines": content.count('\n') + 1,
                    "encoding": encoding if content else "unknown"
                }
            })
        
        # Word Documents
        elif (file_ext in ["docx", "doc"] or 
              mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                           "application/msword"]):
            
            try:
                doc = docx.Document(file_path)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                result.update({
                    "type": "document",
                    "content": content,
                    "metadata": {"paragraphs": len(doc.paragraphs), "pages": len(doc.sections)}
                })
            except Exception as e:
                result["metadata"]["error"] = f"Word processing error: {str(e)}"
        
        # PDFs
        elif (file_ext == "pdf" or 
              mime_type == "application/pdf"):
            
            pdf_text = ""
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        pdf_text += page.extract_text() + "\n"
                
                result.update({
                    "type": "document",
                    "content": pdf_text,
                    "metadata": {"pages": len(pdf_reader.pages)}
                })
            except Exception as e:
                result["metadata"]["error"] = f"PDF processing error: {str(e)}"
        
        # Markdown
        elif (file_ext == "md" or 
              mime_type in ["text/markdown", "text/x-markdown"]):
            
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            html_content = markdown.markdown(md_content)
            result.update({
                "type": "document",
                "content": md_content,
                "html_content": html_content,
                "metadata": {"length": len(md_content)}
            })
        
        # Spreadsheets
        elif (file_ext in ["xlsx", "xls", "csv", "ods"] or 
              mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                           "application/vnd.ms-excel",
                           "text/csv",
                           "application/vnd.oasis.opendocument.spreadsheet"]):
            
            sheets_data = []
            excel_file = None
            
            try:
                if file_ext == "csv" or mime_type == "text/csv":
                    df = pd.read_csv(file_path)
                    sheet_json = {"name": "Sheet1", "headers": df.columns.tolist(), "rows": df.values.tolist(),"metadata": {"formulas": {}, "comments": {}}}
                    sheets_data.append(sheet_json)
                    preview = df.head(5)
                elif file_ext == "ods" or mime_type == "application/vnd.oasis.opendocument.spreadsheet":
                    excel_file = pd.ExcelFile(file_path, engine='odf')
                    for sheet in excel_file.sheet_names:
                        df = excel_file.parse(sheet_name=sheet)
                        sheet_json = {"name": sheet, "headers": df.columns.tolist(), "rows": df.values.tolist(), "metadata": {"formulas": {}, "comments": {}}}
                        sheets_data.append(sheet_json)
                    preview = excel_file.parse(sheet_name=excel_file.sheet_names[0]).head(5)
                else:
                    excel_file = pd.ExcelFile(file_path)
                    for sheet in excel_file.sheet_names:
                        df = excel_file.parse(sheet_name=sheet)
                        sheet_json = {"name": sheet, "headers": df.columns.tolist(), "rows": df.values.tolist(), "metadata": {"formulas": {}, "comments": {}}}
                        sheets_data.append(sheet_json)
                    preview = excel_file.parse(sheet_name=excel_file.sheet_names[0]).head(5)
                
                final_json = {
                    "workbook": filename,
                    "sheets": sheets_data
                }
                
                result.update({
                    "type": "spreadsheet",
                    "content": final_json,
                    "metadata": {
                        "total_sheets": len(sheets_data),
                        "total_rows": sum(len(sheet["rows"]) for sheet in sheets_data),
                        "total_columns": sum(len(sheet["headers"]) for sheet in sheets_data),
                        "sheet_names": [sheet["name"] for sheet in sheets_data]
                    }
                })
                
            except Exception as e:
                result["metadata"]["error"] = f"Spreadsheet processing error: {str(e)}"
        
        # Images
        elif (file_ext in ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp"] or 
              mime_type.startswith("image/")):
            
            try:
                image = Image.open(file_path)
                result.update({
                    "type": "image",
                    "metadata": {
                        "width": image.width,
                        "height": image.height,
                        "format": image.format,
                        "mode": image.mode
                    }
                })
                result["metadata"]["image_data"] = f"Image loaded successfully: {image.size}"
            except Exception as e:
                result["metadata"]["error"] = f"Image processing error: {str(e)}"
        
        # Structured Data (JSON, XML, YAML)
        elif (file_ext in ["json"] or 
              mime_type == "application/json"):
            
            with open(file_path, 'r', encoding='utf-8') as f:
                json_content = json.load(f)
            result.update({
                "type": "structured_data",
                "content": json_content,
                "metadata": {"data_type": type(json_content).__name__}
            })
            preview = json_content if isinstance(json_content, (dict, list)) else {"content": json_content}
        
        elif (file_ext in ["xml"] or 
              mime_type in ["application/xml", "text/xml"]):
            
            try:
                tree = ET.parse(file_path)
                root = tree.getroot()
                # Convert XML to dict for easier handling
                xml_dict = xml_to_dict(root)
                result.update({
                    "type": "structured_data",
                    "content": xml_dict,
                    "metadata": {"root_tag": root.tag}
                })
                preview = xml_dict
            except Exception as e:
                # Fallback to string representation
                with open(file_path, 'r', encoding='utf-8') as f:
                    xml_content = f.read()
                result.update({
                    "type": "structured_data",
                    "content": xml_content,
                    "metadata": {"root_tag": "unknown", "error": str(e)}
                })
                preview = xml_content[:1000] + "..." if len(xml_content) > 1000 else xml_content
        
        elif (file_ext in ["yaml", "yml"] or 
              mime_type in ["application/x-yaml", "text/x-yaml"]):
            
            with open(file_path, 'r', encoding='utf-8') as f:
                yaml_content = yaml.safe_load(f)
            result.update({
                "type": "structured_data",
                "content": yaml_content,
                "metadata": {"data_type": type(yaml_content).__name__}
            })
            preview = yaml_content if isinstance(yaml_content, (dict, list)) else {"content": yaml_content}
        
        
        # Code files
        elif file_ext in ["py", "js", "java", "c", "cpp", "html", "css"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                code_content = f.read()
            result.update({
                "type": "code",
                "content": code_content,
                "metadata": {"lines": code_content.count('\n') + 1, "language": file_ext}
            })
        
        # Presentations (basic text extraction)
        elif file_ext in ["pptx", "odp"]:
            if file_ext == "pptx":
                prs = Presentation(file_path)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                
                result.update({
                    "type": "presentation",
                    "content": content,
                    "metadata": {"slides": len(prs.slides)}
                })
        
        # For unsupported types, try to read as text
        else:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                result.update({
                    "type": "text",
                    "content": content,
                    "metadata": {"length": len(content)}
                })
            except:
                result["metadata"]["error"] = "Unsupported file type"
    
    except Exception as e:
        result["metadata"]["error"] = str(e)
    
    return result, preview


def tabular_to_json(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()
 
    sheets_data = []
 
    if file_type == "csv":
        df = pd.read_csv(uploaded_file)
        sheet_json = {
            "name": "Sheet1",
            "headers": df.columns.tolist(),
            "rows": df.values.tolist(),
            "metadata": {"formulas": {}, "comments": {}}
        }
        sheets_data.append(sheet_json)
 
    else:  # Excel formats
        excel_file = pd.ExcelFile(uploaded_file)
 
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(uploaded_file, sheet_name=sheet)
 
            sheet_json = {
                "name": sheet,
                "headers": df.columns.tolist(),
                "rows": df.values.tolist(),
                "metadata": {"formulas": {}, "comments": {}}
            }
            sheets_data.append(sheet_json)
 
    final_json = {
        "workbook": uploaded_file.name,
        "sheets": sheets_data
    }
 
    return json.dumps(final_json, indent=2)


def parse_markdown_to_json(md_text):
    lines = md_text.split("\n")
    data = {"title": None, "introduction": [], "sections": [], "footnotes": [], "metadata": {}}
   
    current_section = None
    table_headers = []
    table_rows = []
    processing_table = False
    title_found = False
    first_section_found = False
   
    i = 0
    while i < len(lines):
        line = lines[i].strip()
       
        # Skip empty lines and horizontal rules (---)
        if not line or line.startswith("---"):
            i += 1
            continue
       
        # Title (# ...)
        if line.startswith("# "):
            data["title"] = line[2:].strip()
            title_found = True
       
        # Section Heading (## ...)
        elif line.startswith("## "):
            # Mark that we've found the first section
            first_section_found = True
           
            # Finalize previous section
            if current_section:
                data["sections"].append(current_section)
           
            # Start new section
            current_section = {"heading": line[3:].strip()}
            table_headers = []
            table_rows = []
            processing_table = False
       
        # Check for table start
        elif "|" in line and current_section:
            # Look ahead to see if next line is a separator
            if i + 1 < len(lines) and "---" in lines[i + 1] and "|" in lines[i + 1]:
                # This is a table header
                table_headers = [col.strip() for col in line.split("|") if col.strip()]
                processing_table = True
                i += 1  # Skip the separator line
            elif processing_table and table_headers:
                # This is a table row
                row = [col.strip() for col in line.split("|") if col.strip()]
                if row:
                    table_rows.append(row)
            else:
                # Not a table, treat as regular content
                processing_table = False
                if "paragraphs" not in current_section:
                    current_section["paragraphs"] = []
                current_section["paragraphs"].append(line)
       
        # Footnotes
        elif re.match(r"^\d+:", line):
            data["footnotes"].append(line)
       
        # Metadata
        elif ":" in line and any(keyword in line.lower() for keyword in ["author", "created", "date", "version"]):
            parts = line.split(":", 1)
            if len(parts) == 2:
                key = parts[0].strip().lower()
                value = parts[1].strip()
                data["metadata"][key] = value
       
        # Regular content
        else:
            # Handle content after title but before first section (introduction)
            if title_found and not first_section_found:
                # Clean formatting
                content = line
                if content.startswith("- "):
                    content = content[2:].strip()
                elif content.startswith("* "):
                    content = content[2:].strip()
                elif content.startswith("### "):
                    content = content[4:].strip()
               
                if content:
                    data["introduction"].append(content)
           
            # Handle content within sections
            elif current_section:
                # If we were processing a table and hit non-table content, save the table
                if processing_table and table_headers and table_rows:
                    if "tables" not in current_section:
                        current_section["tables"] = []
                    current_section["tables"].append({
                        "headers": table_headers,
                        "rows": table_rows
                    })
                    table_headers = []
                    table_rows = []
                    processing_table = False
               
                # Add to paragraphs
                if "paragraphs" not in current_section:
                    current_section["paragraphs"] = []
               
                # Clean formatting
                content = line
                if content.startswith("- "):
                    content = content[2:].strip()
                elif content.startswith("* "):
                    content = content[2:].strip()
                elif content.startswith("### "):
                    content = content[4:].strip()
               
                if content:
                    current_section["paragraphs"].append(content)
       
        i += 1
   
    # Finalize last section
    if current_section:
        # Save any remaining table
        if processing_table and table_headers and table_rows:
            if "tables" not in current_section:
                current_section["tables"] = []
            current_section["tables"].append({
                "headers": table_headers,
                "rows": table_rows
            })
       
        data["sections"].append(current_section)
   
    return data


######################### xml helper function ###################################
def xml_to_dict(element):
    """Convert XML element to dictionary"""
    result = {}
    for child in element:
        if len(child) == 0:
            result[child.tag] = child.text
        else:
            result[child.tag] = xml_to_dict(child)
    return result

######################### Doc Mime type detection #################################
def detect_file_type(file_content, file_name):
        """
        Detect MIME type of a file using multiple methods for accuracy
        """
        file_ext = Path(file_name).suffix.lower()[1:]
        
        # Method 1: Try using python-magic library (more accurate)
        try:
            mime = magic.Magic(mime=True)
            detected_type = mime.from_buffer(file_content[:2048])  # Read first 2KB for detection
            return detected_type, file_ext
        except:
            pass
        
        # Method 2: Use Streamlit's detected type if available
        if hasattr(file_content, 'type') and file_content.type:
            return file_content.type, file_ext
        
        # Method 3: Use mimetypes library as fallback
        guessed_type, _ = mimetypes.guess_type(file_name)
        if guessed_type:
            return guessed_type, file_ext
        
        # Method 4: Fallback to extension-based detection
        return f"application/{file_ext}" if file_ext else "application/octet-stream", file_ext
    
def categorize_file(mime_type, extension, supported_file_types):
    """
    Categorize file based on MIME type and extension
    """
    for category, info in supported_file_types.items():
        if (mime_type in info.get("mime_types", []) or 
            extension in info.get("extensions", [])):
            return category
    return "Unknown"